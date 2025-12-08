from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Helper function to add a slide with title and bullet points
    def add_slide(title_text, content_list):
        slide_layout = prs.slide_layouts[1] # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set Title
        title = slide.shapes.title
        title.text = title_text
        
        # Set Content
        tf = slide.placeholders[1].text_frame
        tf.clear() # Clear default empty paragraph
        
        for item in content_list:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(20)
            # Check if it's a sub-point (simple heuristic: starts with space or -)
            if item.startswith("  -"):
                p.text = item.replace("  -", "")
                p.level = 1

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "QueueTogether 順便帶"
    subtitle.text = "讓排隊更有價值，實現最後 100 公尺的互助共享\n\n[您的名字/團隊]"

    # Slide 2: 痛點分析
    add_slide("痛點分析：為什麼我們需要這個？", [
        "時間成本高昂：熱門店排隊動輒 30-60 分鐘。",
        "資源浪費：多人分別排隊，不如一人代買。",
        "外送限制：許多名店無外送或溢價過高。",
        "資訊不透明：到了現場才發現人山人海。"
    ])

    # Slide 3: 解決方案
    add_slide("解決方案：共享排隊經濟", [
        "即時媒合：連結「排隊者 (Host)」與「需求者 (Guest)」。",
        "互助共利：Host 賺取跑腿費；Guest 節省時間。",
        "資訊同步：現場實況即時回報。",
        "核心概念：把一個人的等待，轉化為一群人的便利。"
    ])

    # Slide 4: 為什麼選擇 LINE Bot？
    add_slide("產品形式：LINE Bot + LIFF", [
        "最低門檻：無需下載 App，掃碼即用。",
        "原生社交：利用 LINE 群組信任關係，解決面交疑慮。",
        "高滲透率：台灣使用率最高的通訊軟體。",
        "即時通知：Push Message 確保訂單狀態不漏接。"
    ])

    # Slide 5: 使用情境 (User Journey)
    add_slide("使用流程演示", [
        "【Host 發起】：",
        "  - 透過 LIFF 填寫店家與等待時間。",
        "  - 系統生成 Flex Message 卡片傳至群組。",
        "【Guest 跟團】：",
        "  - 點擊卡片直接 +1 下單。",
        "  - 收到 Host「已買到」、「回程中」通知。"
    ])

    # Slide 6: 獨家功能亮點
    add_slide("功能亮點", [
        "智慧截單：設定數量上限或是倒數計時。",
        "實況看板：上傳現場照片，系統預估等待時間。",
        "面交導航：整合 Location Message 精準定位。",
        "信任評分：累積勳章 (如：準時達人)。"
    ])

    # Slide 7: 技術架構
    add_slide("技術架構 (Tech Stack)", [
        "前端：LINE Messaging API + LIFF (HTML/JS)",
        "後端：Python (Flask/FastAPI) on Render",
        "資料庫：PostgreSQL (儲存訂單與狀態)",
        "優勢：輕量開發、快速部署、成本低廉"
    ])

    # Slide 8: 商業模式
    add_slide("商業模式與展望", [
        "初期 (MVP)：累積用戶，建立互助習慣。",
        "中期：小額跑腿費機制、店家導流合作。",
        "願景：成為區域性的「超短距物流」平台。"
    ])

    # Slide 9: 結語
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "感謝聆聽"
    subtitle.text = "讓我們一起把時間花在更美好的事物上\n\n[QR Code 預留區]"

    # Save
    filename = "QueueTogether_Pitch.pptx"
    prs.save(filename)
    print(f"成功生成檔案：{filename}")

if __name__ == "__main__":
    create_presentation()
    
    
    
    
    
    
    
    
line-queue-bot/
├── app.py                # 後端主程式
├── requirements.txt      # 套件依賴清單
└── templates/
    └── liff.html         # 前端 LIFF 頁面



Flask==3.0.0
line-bot-sdk==3.5.0
gunicorn==21.2.0
requests==2.31.0


app.py
import os
import uuid
from datetime import datetime
from flask import Flask, request, abort, render_template, jsonify
from linebot.v3 import WebhookHandler
from linebot.v3.exceptions import InvalidSignatureError
from linebot.v3.messaging import (
    Configuration,
    ApiClient,
    MessagingApi,
    ReplyMessageRequest,
    TextMessage,
    FlexMessage,
    FlexContainer
)
from linebot.v3.webhooks import MessageEvent, TextMessageContent

app = Flask(__name__)

# --- 設定環境變數 ---
CHANNEL_ACCESS_TOKEN = os.getenv('CHANNEL_ACCESS_TOKEN')
CHANNEL_SECRET = os.getenv('CHANNEL_SECRET')
LIFF_ID = os.getenv('LIFF_ID')  # 格式像是: 1234567890-AbCdEfGh

if not all([CHANNEL_ACCESS_TOKEN, CHANNEL_SECRET, LIFF_ID]):
    print("警告: 請檢查環境變數是否設定完整。")

configuration = Configuration(access_token=CHANNEL_ACCESS_TOKEN)
handler = WebhookHandler(CHANNEL_SECRET)

# --- 模擬資料庫 (In-Memory DB) ---
# 在正式環境請改用 PostgreSQL
# 結構: { 'group_id': { 'store':Str, 'host_id':Str, 'orders':List, 'status':Str } }
GROUPS = {}

@app.route("/callback", methods=['POST'])
def callback():
    signature = request.headers['X-Line-Signature']
    body = request.get_data(as_text=True)
    app.logger.info("Request body: " + body)
    try:
        handler.handle(body, signature)
    except InvalidSignatureError:
        abort(400)
    return 'OK'

# --- LIFF 頁面入口 ---
@app.route("/liff")
def liff_entry():
    # 判斷是發起還是跟團
    group_id = request.args.get('group_id')
    return render_template('liff.html', liff_id=LIFF_ID, group_id=group_id)

# --- API: 建立排隊 (Host) ---
@app.route("/api/create_group", methods=['POST'])
def create_group():
    data = request.json
    user_id = data.get('userId')
    user_name = data.get('userName')
    store_name = data.get('storeName')
    
    group_id = str(uuid.uuid4())[:8] # 產生短 ID
    
    GROUPS[group_id] = {
        'id': group_id,
        'store': store_name,
        'host_id': user_id,
        'host_name': user_name,
        'created_at': datetime.now().strftime("%H:%M"),
        'orders': [],
        'status': 'OPEN'
    }
    
    # 建立 Flex Message 卡片
    flex_msg = generate_flex_message(GROUPS[group_id])
    
    # 透過 API 主動推播給 Host (Host 再分享出去) 
    # 註：免費版 LINE Bot 無法主動 Push 給未互動者，
    # 實務上通常建議在 LIFF 用 liff.sendMessages 發送，這邊示範後端回傳邏輯
    
    return jsonify({
        "status": "success", 
        "group_id": group_id,
        "flex_message": flex_msg
    })

# --- API: 加入排隊 (Guest) ---
@app.route("/api/join_group", methods=['POST'])
def join_group():
    data = request.json
    group_id = data.get('groupId')
    user_name = data.get('userName')
    item = data.get('item')
    
    if group_id not in GROUPS:
        return jsonify({"status": "error", "msg": "訂單不存在或已結束"}), 404
        
    GROUPS[group_id]['orders'].append({
        'user': user_name,
        'item': item
    })
    
    return jsonify({"status": "success", "current_count": len(GROUPS[group_id]['orders'])})

# --- 輔助函式: 產生 Flex Message JSON ---
def generate_flex_message(group_data):
    # 這是一個簡單的 Flex Message 結構
    join_url = f"https://liff.line.me/{LIFF_ID}?group_id={group_data['id']}"
    
    bubble = {
        "type": "bubble",
        "hero": {
            "type": "image",
            "url": "https://images.unsplash.com/photo-1561758033-d8f48f85b39e?auto=format&fit=crop&w=600&q=80", # 示意圖
            "size": "full",
            "aspectRatio": "20:13",
            "aspectMode": "cover"
        },
        "body": {
            "type": "box",
            "layout": "vertical",
            "contents": [
                {"type": "text", "text": "排隊揪團中 🍔", "weight": "bold", "size": "xl", "color": "#1DB446"},
                {"type": "text", "text": group_data['store'], "weight": "bold", "size": "xxl", "margin": "md"},
                {"type": "text", "text": f"發起人: {group_data['host_name']}", "size": "sm", "color": "#aaaaaa", "wrap": True},
                {"type": "separator", "margin": "xxl"},
                {"type": "box", "layout": "vertical", "margin": "xxl", "spacing": "sm", "contents": [
                    {"type": "box", "layout": "baseline", "spacing": "sm", "contents": [
                        {"type": "text", "text": "時間", "color": "#aaaaaa", "size": "sm", "flex": 1},
                        {"type": "text", "text": group_data['created_at'], "wrap": True, "color": "#666666", "size": "sm", "flex": 5}
                    ]}
                ]}
            ]
        },
        "footer": {
            "type": "box",
            "layout": "vertical",
            "spacing": "sm",
            "contents": [
                {
                    "type": "button",
                    "style": "primary",
                    "height": "sm",
                    "action": {
                        "type": "uri",
                        "label": "我要跟團 +1",
                        "uri": join_url
                    },
                    "color": "#00b900"
                }
            ]
        }
    }
    return bubble

if __name__ == "__main__":
    app.run(debug=True)




html


<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>QueueTogether</title>
    <script src="https://static.line-scdn.net/liff/edge/2/sdk.js"></script>
    <style>
        body { font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif; padding: 20px; background-color: #f5f5f5; }
        .card { background: white; padding: 20px; border-radius: 10px; box-shadow: 0 2px 5px rgba(0,0,0,0.1); }
        input, select, button { width: 100%; padding: 12px; margin-top: 10px; border-radius: 5px; border: 1px solid #ddd; box-sizing: border-box; }
        button { background-color: #00b900; color: white; border: none; font-weight: bold; cursor: pointer; }
        button:disabled { background-color: #ccc; }
        h2 { color: #333; margin-top: 0; }
        .hidden { display: none; }
    </style>
</head>
<body>
    <div class="card">
        <div id="loading">載入中...</div>

        <div id="host-view" class="hidden">
            <h2>🍔 發起排隊</h2>
            <input type="text" id="storeName" placeholder="輸入店家名稱 (如: 50嵐)">
            <button id="createBtn" onclick="createGroup()">建立揪團卡片</button>
        </div>

        <div id="guest-view" class="hidden">
            <h2>📝 我要跟團</h2>
            <p id="groupInfo">正在加入團購...</p>
            <input type="text" id="orderItem" placeholder="你想吃/喝什麼？(如: 珍奶半糖)">
            <button id="joinBtn" onclick="joinGroup()">送出訂單</button>
        </div>
    </div>

    <script>
        // 從後端傳來的變數
        const LIFF_ID = "{{ liff_id }}"; 
        const GROUP_ID = "{{ group_id }}"; // 如果是 None 則為空字串

        async function main() {
            await liff.init({ liffId: LIFF_ID });
            
            if (!liff.isLoggedIn()) {
                liff.login();
                return;
            }

            const profile = await liff.getProfile();
            window.currentUser = profile;
            document.getElementById('loading').style.display = 'none';

            // 判斷模式
            if (GROUP_ID && GROUP_ID !== 'None') {
                // 跟團模式
                document.getElementById('guest-view').classList.remove('hidden');
                document.getElementById('groupInfo').innerText = `加入訂單 ID: ${GROUP_ID}`;
            } else {
                // 發起模式
                document.getElementById('host-view').classList.remove('hidden');
            }
        }

        // Host: 建立群組並發送卡片
        async function createGroup() {
            const store = document.getElementById('storeName').value;
            if (!store) return alert('請輸入店家名稱');

            const res = await fetch('/api/create_group', {
                method: 'POST',
                headers: {'Content-Type': 'application/json'},
                body: JSON.stringify({
                    userId: window.currentUser.userId,
                    userName: window.currentUser.displayName,
                    storeName: store
                })
            });
            
            const data = await res.json();
            
            if (data.status === 'success') {
                // 使用 LIFF API 直接在聊天室發送卡片
                await liff.sendMessages([{
                    type: "flex",
                    altText: "有人發起排隊囉！",
                    contents: data.flex_message
                }]);
                liff.closeWindow();
            }
        }

        // Guest: 加入訂單
        async function joinGroup() {
            const item = document.getElementById('orderItem').value;
            if (!item) return alert('請輸入餐點');

            const res = await fetch('/api/join_group', {
                method: 'POST',
                headers: {'Content-Type': 'application/json'},
                body: JSON.stringify({
                    groupId: GROUP_ID,
                    userName: window.currentUser.displayName,
                    item: item
                })
            });

            const data = await res.json();
            if (data.status === 'success') {
                alert('成功加入！目前人數: ' + data.current_count);
                liff.closeWindow();
            } else {
                alert('錯誤: ' + data.msg);
            }
        }

        main();
    </script>
</body>
</html>


